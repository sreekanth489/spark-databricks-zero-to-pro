"""
Generate PowerPoint presentation for Lakeflow: Connect, Spark Declarative Pipelines, and Jobs.
Run: python3 resources/generate_lakeflow_pptx.py
Output: resources/lakeflow-connect-sdp-jobs.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# -- Color palette --
BG_DARK = RGBColor(0x1B, 0x1F, 0x23)
BG_SECTION = RGBColor(0x23, 0x29, 0x2E)
ACCENT_RED = RGBColor(0xE8, 0x3E, 0x3E)
ACCENT_BLUE = RGBColor(0x42, 0x9E, 0xF5)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
ACCENT_PURPLE = RGBColor(0xAB, 0x47, 0xBC)
ACCENT_TEAL = RGBColor(0x00, 0x96, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBE, 0xC5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def text_box(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def bullets(tf, items, size=16, color=LIGHT_GRAY):
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_before = Pt(4)


def code_box(slide, left, top, w, h, code, font_size=14):
    shape = add_rect(slide, left, top, w, h, RGBColor(0x12, 0x15, 0x18), ACCENT_BLUE)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(code.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN
        p.font.name = "Courier New"
    return shape


def card(slide, left, top, w, h, title, desc, title_color, border_color):
    box = add_rect(slide, left, top, w, h, BG_SECTION, border_color)
    tf = set_text(box, title, size=20, color=title_color, bold=True, align=PP_ALIGN.CENTER)
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)


def table_slide(slide, headers, rows, top_start, col_widths=None, header_color=ACCENT_BLUE):
    n_cols = len(headers)
    if col_widths is None:
        total = 12.0
        cw = total / n_cols
        col_widths = [Inches(cw)] * n_cols
    start_left = Inches(0.5)

    for i, header in enumerate(headers):
        left = start_left + sum(col_widths[:i])
        box = add_rect(slide, left, top_start, col_widths[i], Inches(0.5), header_color)
        set_text(box, header, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    for r, row in enumerate(rows):
        row_top = top_start + Inches(0.55 * (r + 1))
        for c, cell in enumerate(row):
            left = start_left + sum(col_widths[:c])
            box = add_rect(slide, left, row_top, col_widths[c], Inches(0.5), BG_SECTION)
            set_text(box, cell, size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)
text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "Lakeflow", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
         "Connect  |  Spark Declarative Pipelines  |  Jobs", size=30, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
         "Ingest, Transform, and Orchestrate -- The Complete Data Engineering Platform", size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(5.0), Inches(2.3), Inches(0.06), ACCENT_RED)
text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.5),
         "Spark Databricks Zero-to-Pro  |  Day 22, 23 & 24", size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Lakeflow Platform Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_BLUE)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "The Lakeflow Platform", size=36, color=WHITE, bold=True)

platform_code = """                     LAKEFLOW PLATFORM
  ┌──────────────┐      ┌──────────────────────┐      ┌──────────────┐
  │   CONNECT    │ ───> │  SPARK DECLARATIVE    │ ───> │     JOBS     │
  │   (Ingest)   │      │     PIPELINES         │      │ (Orchestrate)│
  │   Day 22     │      │    (Transform)        │      │   Day 24     │
  │              │      │     Day 23            │      │              │
  └──────────────┘      └──────────────────────┘      └──────────────┘

  External Sources       Bronze -> Silver              Scheduling
  -> Bronze Tables       -> Gold Tables                & Monitoring"""
code_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.5), platform_code, font_size=15)

cards_data = [
    ("CONNECT", "Get the data IN\nDatabases, SaaS, Files, APIs", ACCENT_ORANGE, ACCENT_ORANGE),
    ("SDP", "Make the data RIGHT\nBronze -> Silver -> Gold", ACCENT_GREEN, ACCENT_GREEN),
    ("JOBS", "Run it RELIABLY\nSchedule, Monitor, Retry", ACCENT_PURPLE, ACCENT_PURPLE),
]
for i, (title, desc, tc, bc) in enumerate(cards_data):
    card(slide, Inches(0.5 + i * 4.2), Inches(5.2), Inches(3.8), Inches(2), title, desc, tc, bc)

# ============================================================
# SLIDE 3: Section Divider -- Connect
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_SECTION)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_ORANGE)
text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
         "Part 1: Lakeflow Connect", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "Bring External Data INTO the Lakehouse", size=24, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.06), ACCENT_ORANGE)

# ============================================================
# SLIDE 4: Connect Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_ORANGE)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Lakeflow Connect Architecture", size=36, color=WHITE, bold=True)

arch_code = """External Sources              Lakeflow Connect                  Lakehouse
===============        ============================       ==================

Cloud Storage ───────> Standard Connectors           ───> Unity Catalog
(S3/ADLS/GCS)          * Auto Loader (cloudFiles)           ┌───────────┐
                       * Batch (spark.read)                  │  Bronze   │
Kafka Topics  ───────> * Kafka Connector              ───>  │  Tables   │
                       * Streaming (readStream)              └───────────┘
                                                             ┌───────────┐
Databases     ───────> Managed Connectors              ───> │ Streaming │
(PostgreSQL,           * No-code UI setup                    │  Tables   │
 MySQL, Oracle)        * CDC-based incremental               └───────────┘
                       * Serverless compute                  ┌───────────┐
SaaS Apps     ───────> Managed Connectors              ───> │  Volumes  │
(Salesforce,           * Schema auto-inference               └───────────┘
 Workday)              * Auto-scaling

Local Files   ───────> Manual File Upload              ───> Volume or Table"""
code_box(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.8), arch_code, font_size=13)

# ============================================================
# SLIDE 5: Three Ingestion Types
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_ORANGE)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Three Ingestion Types", size=36, color=WHITE, bold=True)

ingestion_types = [
    ("Manual File Upload", ACCENT_BLUE, [
        "Drag & drop via Databricks UI",
        "Best for: one-time loads, <2 GB",
        "Schema auto-inferred",
        "Stored in Volumes or Tables",
    ]),
    ("Standard Connectors", ACCENT_GREEN, [
        "Code-based (PySpark / SQL)",
        "Auto Loader, JDBC, Kafka",
        "Batch, Incremental, Streaming",
        "Full control over ingestion logic",
    ]),
    ("Managed Connectors", ACCENT_PURPLE, [
        "No-code, serverless, CDC-based",
        "Databases + SaaS apps",
        "Auto schema inference + evolution",
        "Built-in monitoring dashboards",
    ]),
]

for i, (title, color, items) in enumerate(ingestion_types):
    left = Inches(0.3 + i * 4.3)
    box = add_rect(slide, left, Inches(1.3), Inches(4.0), Inches(5.5), BG_SECTION, color)
    tf = set_text(box, title, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    tf.add_paragraph()
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(8)

# ============================================================
# SLIDE 6: Standard Connectors Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_GREEN)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Standard Connectors: Ingestion Modes", size=36, color=WHITE, bold=True)

table_slide(slide,
            ["Mode", "Description", "Latency", "Use Case"],
            [
                ["Batch", "Full load every run", "Minutes-Hours", "Small tables, snapshots"],
                ["Incremental Batch", "Only new/changed rows", "Minutes", "Append-only logs"],
                ["Streaming", "Continuous real-time", "Seconds", "Kafka, clickstream, IoT"],
            ],
            Inches(1.3),
            [Inches(2.5), Inches(3.5), Inches(2.5), Inches(3.5)])

std_code = """# Auto Loader (most common standard connector)
spark.readStream
    .format("cloudFiles")
    .option("cloudFiles.format", "json")
    .option("cloudFiles.schemaLocation", "s3://bucket/schemas/")
    .load("s3://bucket/raw/clickstream/")
    .writeStream
    .option("checkpointLocation", "s3://bucket/checkpoints/")
    .toTable("ecommerce.bronze.clickstream")"""
code_box(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(3.4), std_code, font_size=14)

# ============================================================
# SLIDE 7: Managed Connectors
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_PURPLE)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Managed Connectors", size=36, color=WHITE, bold=True)

text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "No-code, serverless, CDC-based ingestion from databases and SaaS apps", size=18, color=ACCENT_BLUE)

features = [
    ("No-code setup", "Configure entirely through Databricks UI"),
    ("Serverless compute", "No clusters to manage; auto-scaling"),
    ("CDC-based", "Uses database change logs for incremental reads"),
    ("Schema evolution", "Handles new columns and type changes"),
    ("Unity Catalog", "All data governed from ingestion"),
    ("Built-in monitoring", "Health, freshness, and error dashboards"),
]
for i, (title, desc) in enumerate(features):
    row = i // 3
    col = i % 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.9 + row * 1.5)
    box = add_rect(slide, left, top, Inches(3.8), Inches(1.2), BG_SECTION, ACCENT_PURPLE)
    tf = set_text(box, title, size=16, color=ACCENT_PURPLE, bold=True, align=PP_ALIGN.CENTER)
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

mc_code = """-- SQL: Create a managed connection + streaming table
CREATE CONNECTION my_postgres TYPE postgresql
OPTIONS (host 'db.example.com', port '5432',
         user secret('jdbc', 'user'), password secret('jdbc', 'pwd'));

CREATE STREAMING TABLE ecommerce.bronze.pg_customers
AS SELECT * FROM STREAM read_changefeed('my_postgres', 'public.customers');"""
code_box(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2), mc_code, font_size=14)

# ============================================================
# SLIDE 8: Comparing Ingestion Methods
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_ORANGE)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Comparing Ingestion Methods", size=36, color=WHITE, bold=True)

table_slide(slide,
            ["Criteria", "Manual Upload", "Standard Connectors", "Managed Connectors"],
            [
                ["Setup", "None (UI drag-drop)", "Code required", "No-code (UI/SQL)"],
                ["Compute", "Workspace", "User-managed clusters", "Serverless"],
                ["Latency", "Manual trigger", "Batch to streaming", "Near real-time (CDC)"],
                ["Scale", "Small (<2 GB)", "Any size", "Auto-scaling"],
                ["Schema Mgmt", "Auto-inferred", "Code-managed", "Auto + evolution"],
                ["Sources", "Local files", "Any Spark source", "Databases + SaaS"],
                ["Monitoring", "None", "Custom", "Built-in dashboards"],
                ["Best For", "Ad-hoc loads", "Custom pipelines", "Enterprise integration"],
            ],
            Inches(1.2),
            [Inches(2.0), Inches(2.8), Inches(3.4), Inches(3.8)])

# ============================================================
# SLIDE 9: Section Divider -- SDP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_SECTION)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_GREEN)
text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
         "Part 2: Spark Declarative Pipelines", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "Declare WHAT Your Data Should Look Like, Not HOW to Build It", size=22, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.06), ACCENT_GREEN)

# ============================================================
# SLIDE 10: Evolution of Spark Pipelines
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_GREEN)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "The Three Eras of Spark Pipelines", size=36, color=WHITE, bold=True)

table_slide(slide,
            ["Era", "Approach", "You Manage", "Framework Manages"],
            [
                ["Traditional Spark (2014+)", "Imperative", "Everything", "Query optimization"],
                ["Delta Live Tables (2021+)", "Declarative", "Pipeline config", "DAG, checkpoints, CDC"],
                ["Lakeflow SDP (2025+)", "Declarative + Platform", "Business logic only", "Everything + orchestration"],
            ],
            Inches(1.2),
            [Inches(3.5), Inches(2.5), Inches(3.0), Inches(3.0)])

analogy_code = """Traditional Spark = Cooking manually         (you manage everything)
DLT              = Smart cooking assistant   (handles timing & quality)
Lakeflow SDP     = Full restaurant kitchen   (manages the entire operation)"""
code_box(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(1.5), analogy_code, font_size=16)

# What each solved
items = [
    ("Checkpoint management", "Manual", "Auto", "Auto"),
    ("Data quality tracking", "None", "Expectations", "Expectations"),
    ("Pipeline DAG", "None", "Auto", "Auto"),
    ("CDC", "Manual MERGE", "apply_changes", "Auto CDC Flow"),
    ("Open source", "Yes", "No (Databricks)", "Yes (Apache Spark)"),
]
top = Inches(5.2)
for header in ["Capability", "Traditional", "DLT", "SDP"]:
    pass  # already shown in table above

table_slide(slide,
            ["Capability", "Traditional", "DLT", "Lakeflow SDP"],
            items,
            Inches(5.1),
            [Inches(3.5), Inches(2.5), Inches(3.0), Inches(3.0)])

# ============================================================
# SLIDE 11: Core Abstractions -- Dataset Types
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_GREEN)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Core Abstractions: Dataset Types", size=36, color=WHITE, bold=True)

dataset_types = [
    ("Streaming Table", ACCENT_BLUE, [
        "Query: Streaming (append-only)",
        "Refresh: Incremental",
        "Storage: Delta table",
        "Best for: Ingestion / Bronze",
        "Python: @dp.table",
        "SQL: CREATE STREAMING TABLE",
    ]),
    ("Materialized View", ACCENT_GREEN, [
        "Query: Batch (full recompute)",
        "Refresh: Complete",
        "Storage: Delta table",
        "Best for: Transforms / Gold",
        "Python: @dp.materialized_view",
        "SQL: CREATE MATERIALIZED VIEW",
    ]),
    ("View", ACCENT_ORANGE, [
        "Query: Batch (on-read)",
        "Refresh: No storage",
        "Storage: None",
        "Best for: Reusable logic",
        "Python: @dp.view",
        "SQL: CREATE VIEW / TEMP VIEW",
    ]),
]

for i, (title, color, items) in enumerate(dataset_types):
    left = Inches(0.3 + i * 4.3)
    box = add_rect(slide, left, Inches(1.3), Inches(4.0), Inches(5.8), BG_SECTION, color)
    tf = set_text(box, title, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    tf.add_paragraph()
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(8)

# ============================================================
# SLIDE 12: SDP Code Examples
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_GREEN)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "SDP: Python & SQL Syntax", size=36, color=WHITE, bold=True)

py_code = """# Python: Streaming Table (Bronze)
from pyspark import pipelines as dp

@dp.table(name="bronze_orders",
          comment="Raw orders via Auto Loader")
def bronze_orders():
    return spark.readStream.format("cloudFiles") \\
        .option("cloudFiles.format", "csv") \\
        .load("s3://bucket/orders")

# Python: Materialized View (Silver/Gold)
@dp.materialized_view(name="silver_stores")
def silver_stores():
    return spark.read.table("bronze_stores") \\
        .select("store_id", "store_name", "city")"""
code_box(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(5.2), py_code, font_size=13)

sql_code = """-- SQL: Streaming Table (Bronze)
CREATE OR REFRESH STREAMING TABLE bronze_orders
COMMENT 'Raw orders via Auto Loader'
AS SELECT * FROM STREAM read_files(
  's3://bucket/orders',
  format => 'csv', header => true
);

-- SQL: Materialized View (Silver/Gold)
CREATE OR REFRESH MATERIALIZED VIEW
  silver_stores
AS SELECT store_id, store_name, city
   FROM bronze_stores;"""
code_box(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.2), sql_code, font_size=13)

# ============================================================
# SLIDE 13: Data Quality with Expectations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_TEAL)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Data Quality: Expectations", size=36, color=WHITE, bold=True)

# Enforcement levels
levels = [
    ("expect", ACCENT_BLUE, "Bad records KEPT\nPipeline continues\nUse: Monitor quality"),
    ("expect_or_drop", ACCENT_ORANGE, "Bad records DROPPED\nPipeline continues\nUse: Filter bad data"),
    ("expect_or_fail", ACCENT_RED, "Pipeline FAILS\nImmediate halt\nUse: Hard constraints"),
]
for i, (title, color, desc) in enumerate(levels):
    card(slide, Inches(0.3 + i * 4.3), Inches(1.3), Inches(4.0), Inches(2.2), title, desc, color, color)

py_exp = """# Python
@dp.table(name="silver_orders")
@dp.expect("valid_amount", "order_amount > 0")
@dp.expect_or_drop("valid_rating", "customer_rating BETWEEN 1 AND 5")
@dp.expect_or_fail("valid_id", "order_id IS NOT NULL")
def silver_orders():
    return spark.read.table("bronze_orders")"""
code_box(slide, Inches(0.3), Inches(3.8), Inches(6.2), Inches(3.2), py_exp, font_size=13)

sql_exp = """-- SQL
CREATE OR REFRESH STREAMING TABLE silver_orders (
  CONSTRAINT valid_amount
    EXPECT (order_amount > 0),
  CONSTRAINT valid_rating
    EXPECT (customer_rating BETWEEN 1 AND 5)
    ON VIOLATION DROP ROW,
  CONSTRAINT valid_id
    EXPECT (order_id IS NOT NULL)
    ON VIOLATION FAIL UPDATE
) AS SELECT * FROM STREAM bronze_orders;"""
code_box(slide, Inches(6.8), Inches(3.8), Inches(6.2), Inches(3.2), sql_exp, font_size=13)

# ============================================================
# SLIDE 14: Auto CDC Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_GREEN)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Change Data Capture: Auto CDC Flow", size=36, color=WHITE, bold=True)

cdc_py = """# Python: Auto CDC Flow
import pyspark.sql.functions as F
from pyspark import pipelines as dp

dp.create_streaming_table(name="silver_orders",
                          comment="Orders with CDC applied")

dp.create_auto_cdc_flow(
    name="silver_orders_cdc",
    target="silver_orders",
    source="bronze_orders_staging",
    keys=["order_id"],
    sequence_by="updated_at",
    stored_as_scd_type=1,
    apply_as_deletes=F.expr("operation = 'DELETE'")
)"""
code_box(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(4.2), cdc_py, font_size=13)

cdc_sql = """-- SQL: Auto CDC Flow
CREATE OR REFRESH STREAMING TABLE silver_orders;

APPLY CHANGES INTO silver_orders
FROM bronze_orders_staging
KEYS (order_id)
SEQUENCE BY updated_at
STORED AS SCD TYPE 1;

-- SCD Type 2: tracks full history
-- STORED AS SCD TYPE 2
-- Adds: __start_at, __end_at, __is_current"""
code_box(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(4.2), cdc_sql, font_size=13)

# SCD comparison
table_slide(slide,
            ["Feature", "SCD Type 1", "SCD Type 2"],
            [
                ["History", "No -- overwrites", "Yes -- full history"],
                ["Extra Columns", "None", "__start_at, __end_at, __is_current"],
                ["Storage", "Less", "More (row per version)"],
                ["Use Case", "Current state only", "Audit trail, historical"],
            ],
            Inches(5.7),
            [Inches(3.0), Inches(4.5), Inches(4.5)])

# ============================================================
# SLIDE 15: Pipeline DAG & Benefits
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_GREEN)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Pipeline DAG & Key Benefits", size=36, color=WHITE, bold=True)

dag_code = """  S3: orders/           S3: stores/
      |                     |
      v                     v
+--------------+    +--------------+
| bronze.orders|    | bronze.stores|
| (Streaming   |    | (Materialized|
|  Table)      |    |  View)       |
+--------------+    +--------------+
      |                     |
      v                     v
+--------------+    +--------------+    +--------------+
| silver.orders|    | silver.stores|    | silver.      |
| (CDC Flow)   |    | (Mat. View)  |    |  calendar    |
+--------------+    +--------------+    +--------------+
      |                  |    |               |
      +--------+---------+    +-------+-------+
               |                      |
               v                      v
       +---------------+    +------------------+
       | gold.fact_    |    | gold.regional_   |
       |   orders      |    |   views          |
       +---------------+    +------------------+"""
code_box(slide, Inches(0.3), Inches(1.2), Inches(6.5), Inches(5.8), dag_code, font_size=11)

benefits = [
    ("Pre-Validation", "Validates entire graph\nbefore execution"),
    ("Auto Checkpoints", "No checkpointLocation\nneeded ever"),
    ("Auto Parallelism", "Independent steps run\nsimultaneously"),
    ("Dependency Resolution", "Order discovered from\nyour queries"),
    ("Efficient Retries", "Only failed steps\nare retried"),
    ("Event Logging", "Full audit trail\nof pipeline runs"),
]
for i, (title, desc) in enumerate(benefits):
    row = i // 2
    col = i % 2
    left = Inches(7.2 + col * 3.0)
    top = Inches(1.3 + row * 2.0)
    box = add_rect(slide, left, top, Inches(2.7), Inches(1.7), BG_SECTION, ACCENT_GREEN)
    tf = set_text(box, title, size=15, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)

# ============================================================
# SLIDE 16: Pipeline Modes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_GREEN)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Pipeline Configuration", size=36, color=WHITE, bold=True)

text_box(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.5),
         "Pipeline Modes", size=22, color=ACCENT_BLUE, bold=True)
table_slide(slide,
            ["Mode", "Behavior", "Best For"],
            [
                ["Triggered", "Runs once, then stops", "Scheduled batch jobs"],
                ["Continuous", "Keeps running", "Low-latency streaming"],
            ],
            Inches(1.7),
            [Inches(2.0), Inches(4.0), Inches(4.0)])

text_box(slide, Inches(0.8), Inches(3.2), Inches(5), Inches(0.5),
         "Development vs Production", size=22, color=ACCENT_BLUE, bold=True)
table_slide(slide,
            ["Setting", "Development Mode", "Production Mode"],
            [
                ["Cluster", "Reused across runs", "New cluster per run"],
                ["Retries", "Disabled (fail fast)", "Enabled (resilient)"],
                ["Purpose", "Fast iteration", "Reliable production"],
            ],
            Inches(3.7),
            [Inches(2.0), Inches(4.0), Inches(4.0)])

text_box(slide, Inches(0.8), Inches(5.5), Inches(5), Inches(0.5),
         "Cluster Modes", size=22, color=ACCENT_BLUE, bold=True)
table_slide(slide,
            ["Mode", "Description"],
            [
                ["Fixed Size", "Static workers; predictable cost"],
                ["Enhanced Autoscaling", "SDP-optimized; recommended for production"],
                ["Legacy Autoscaling", "Standard Spark; compatibility only"],
            ],
            Inches(6.0),
            [Inches(3.5), Inches(6.5)])

# ============================================================
# SLIDE 17: Section Divider -- Jobs
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_SECTION)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_PURPLE)
text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
         "Part 3: Lakeflow Jobs", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "Orchestrate the Entire Workflow: Schedule, Monitor, Retry", size=22, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.06), ACCENT_PURPLE)

# ============================================================
# SLIDE 18: Multi-Task Job Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_PURPLE)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Multi-Task Job Architecture", size=36, color=WHITE, bold=True)

text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "Supported Task Types", size=20, color=ACCENT_BLUE, bold=True)

task_types = [
    ["Notebook", "SDP Pipeline", "SQL", "Python Script", "JAR"],
    ["dbt", "Spark Submit", "Run Job", "If/Else", "For Each"],
]
for r, row in enumerate(task_types):
    for c, task in enumerate(row):
        left = Inches(0.5 + c * 2.5)
        top = Inches(1.9 + r * 0.65)
        box = add_rect(slide, left, top, Inches(2.3), Inches(0.55), BG_SECTION, ACCENT_PURPLE)
        set_text(box, task, size=14, color=LIGHT_GRAY, bold=True, align=PP_ALIGN.CENTER)

# DAG Patterns
text_box(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.5),
         "DAG Patterns", size=20, color=ACCENT_BLUE, bold=True)

patterns_code = """Pattern 1: Linear            Pattern 2: Fan-Out / Fan-In
A --> B --> C                         +-> B (US) --+
                              A ------+-> C (EU) --+-> E (aggregate)
                                      +-> D (APAC)-+

Pattern 3: Conditional        Pattern 4: End-to-End Lakeflow
         +-- success --> B    Task 1: Connect --> Task 2: SDP Pipeline
A -------+                                             |
         +-- failure --> C            +----------------+-----------+
                                      |                           |
                                Task 3: SQL Report        Task 4: Notify"""
code_box(slide, Inches(0.3), Inches(4.2), Inches(12.7), Inches(3), patterns_code, font_size=13)

# ============================================================
# SLIDE 19: Trigger Modes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_PURPLE)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Trigger Modes", size=36, color=WHITE, bold=True)

table_slide(slide,
            ["Mode", "When It Runs", "Best For", "Configuration"],
            [
                ["Manual", "On-demand (UI/API/CLI)", "Ad-hoc, testing", "No trigger needed"],
                ["Scheduled", "Cron-based schedule", "Nightly ETL, hourly", "Cron + timezone"],
                ["Continuous", "After previous run ends", "Near-real-time", "continuous flag"],
                ["File Arrival", "New files in storage", "Event-driven ingestion", "S3/ADLS path + wait"],
            ],
            Inches(1.2),
            [Inches(2.0), Inches(3.0), Inches(3.5), Inches(3.5)])

trigger_code = """# Cron examples
0 0 2 * * ?      # Every day at 2 AM UTC
0 0 * * * 1-5    # Every hour on weekdays
0 */15 * * * ?   # Every 15 minutes

# File Arrival Trigger:
Path:       s3://ecommerce-lakehouse/raw/orders/
Min files:  1
Wait time:  60 seconds (debounce period)

# Continuous Trigger:
Run 1 starts -> Run 1 ends -> Run 2 starts -> Run 2 ends -> ..."""
code_box(slide, Inches(0.5), Inches(3.9), Inches(12.3), Inches(3.2), trigger_code, font_size=14)

# ============================================================
# SLIDE 20: Repair Runs
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Repair Runs: Re-run Only Failed Tasks", size=36, color=WHITE, bold=True)

repair_code = """Original Run:                              Repair Run:
  Task A (ingest)     -- SUCCESS (30 min)    Task A -- SKIPPED (reuses result)
  Task B (transform)  -- SUCCESS (45 min)    Task B -- SKIPPED (reuses result)
  Task C (aggregate)  -- FAILED  (10 min)    Task C -- RE-RUN  (after fix)
  Task D (report)     -- SKIPPED             Task D -- RE-RUN  (depends on C)

  Total: 85 min wasted                       Total: only ~15 min to repair!"""
code_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.8), repair_code, font_size=15)

repair_benefits = [
    "Saves significant compute time and cost",
    "Only failed tasks and their downstream dependencies are re-executed",
    "Successful tasks reuse their previous results",
    "No need to re-run the entire pipeline from scratch",
    "Available via UI (Repair Run button) or API",
]
tf = text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(2.5), "", size=17)
for item in repair_benefits:
    p = tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(17)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

# ============================================================
# SLIDE 21: Cluster Strategies & RBAC
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_PURPLE)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Cluster Strategies & RBAC", size=36, color=WHITE, bold=True)

text_box(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
         "Cluster Strategies", size=20, color=ACCENT_BLUE, bold=True)
table_slide(slide,
            ["Strategy", "Cost", "Use Case"],
            [
                ["Job cluster", "Low", "Production batch jobs"],
                ["Shared job cluster", "Medium", "Multi-task same compute"],
                ["All-purpose cluster", "High", "Development / debugging"],
                ["Serverless", "Variable", "No management overhead"],
            ],
            Inches(1.6),
            [Inches(3.0), Inches(1.5), Inches(5.5)])

text_box(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.5),
         "RBAC Permissions", size=20, color=ACCENT_BLUE, bold=True)
table_slide(slide,
            ["Permission", "Capabilities"],
            [
                ["Can View", "See config and run history"],
                ["Can Manage Run", "View + trigger/cancel runs"],
                ["Can Manage", "Full control: edit, delete, permissions"],
                ["Is Owner", "Manage + transfer ownership"],
            ],
            Inches(4.8),
            [Inches(3.0), Inches(7.0)])

# ============================================================
# SLIDE 22: Databricks Asset Bundles
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_TEAL)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Databricks Asset Bundles (DAB)", size=36, color=WHITE, bold=True)

text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "CI/CD for Databricks: Jobs defined as code, deployed across environments", size=18, color=ACCENT_BLUE)

dab_code = """my-project/
+-- databricks.yml        # Main configuration (YAML)
+-- resources/
|   +-- job_config.yml    # Job definitions
+-- src/
|   +-- ingest_orders.py
|   +-- transform.py
+-- tests/
    +-- test_pipeline.py

# Key CLI Commands:
databricks bundle init                   # Initialize project
databricks bundle validate               # Validate config
databricks bundle deploy --target dev    # Deploy to dev
databricks bundle deploy --target prod   # Deploy to prod
databricks bundle run my_job --target dev  # Run a job"""
code_box(slide, Inches(0.3), Inches(1.8), Inches(6.2), Inches(4.5), dab_code, font_size=13)

dab_yaml = """# databricks.yml -- Environment Targets
targets:
  dev:
    workspace:
      host: https://dbc-abc123.cloud.databricks.com
    default: true
  staging:
    workspace:
      host: https://dbc-def456.cloud.databricks.com
  prod:
    workspace:
      host: https://dbc-ghi789.cloud.databricks.com
    run_as:
      service_principal_name: "prod-sp" """
code_box(slide, Inches(6.8), Inches(1.8), Inches(6.2), Inches(4.5), dab_yaml, font_size=13)

comp_items = [
    ("Without DAB", "Manual UI job creation, copy/paste between envs, no version control"),
    ("With DAB", "Jobs as code, automated deployment, full git history, pre-validated"),
]
for i, (title, desc) in enumerate(comp_items):
    color = ACCENT_RED if i == 0 else ACCENT_GREEN
    left = Inches(0.3 + i * 6.5)
    box = add_rect(slide, left, Inches(6.5), Inches(6.2), Inches(0.7), BG_SECTION, color)
    tf = set_text(box, f"{title}: ", size=14, color=color, bold=True)
    run = tf.paragraphs[0].add_run()
    run.text = desc
    run.font.size = Pt(13)
    run.font.color.rgb = LIGHT_GRAY

# ============================================================
# SLIDE 23: End-to-End Lakeflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "End-to-End Lakeflow Workflow", size=36, color=WHITE, bold=True)

e2e_code = """                        LAKEFLOW JOBS (Orchestration Layer)
  ┌─────────────────────────────────────────────────────────────────────────┐
  │                                                                         │
  │   Task 1: CONNECT               Task 2: SDP PIPELINE                   │
  │   ┌───────────────────┐         ┌────────────────────────────────┐     │
  │   │ Auto Loader       │         │ bronze_orders (Streaming Table)│     │
  │   │ S3 -> Bronze      │ ──────> │ silver_orders (CDC + Expects)  │     │
  │   │                   │         │ gold_summary  (Mat. View)      │     │
  │   └───────────────────┘         └────────────────────────────────┘     │
  │                                          │                              │
  │                              ┌───────────┴──────────┐                  │
  │                              │                      │                  │
  │                    Task 3: SQL Report      Task 4: Python Notify       │
  │                    ┌────────────────┐      ┌──────────────────┐       │
  │                    │ Daily revenue  │      │ Slack + email    │       │
  │                    │ aggregation    │      │ notification     │       │
  │                    └────────────────┘      └──────────────────┘       │
  │                                                                         │
  │   Trigger: Scheduled (daily 2 AM)  |  Retry: 3 attempts per task      │
  │   Cluster: Job cluster (auto-term) |  Alerts: on_failure -> Slack     │
  └─────────────────────────────────────────────────────────────────────────┘"""
code_box(slide, Inches(0.2), Inches(1.2), Inches(12.9), Inches(5.8), e2e_code, font_size=12)

# ============================================================
# SLIDE 24: Certification Tips
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Certification Tips", size=36, color=WHITE, bold=True)

sections = [
    ("Lakeflow Connect", ACCENT_ORANGE, [
        "Know the 3 ingestion types: Manual, Standard, Managed",
        "Auto Loader = standard connector for cloud files",
        "Managed Connectors: serverless, no-code, CDC-based",
        "SQL: CREATE CONNECTION + CREATE STREAMING TABLE",
    ]),
    ("Spark Declarative Pipelines", ACCENT_GREEN, [
        "3 enforcement levels: expect, expect_or_drop, expect_or_fail",
        "Streaming Table vs Materialized View: when to use each",
        "Auto CDC: KEYS, SEQUENCE BY, SCD Type 1 vs 2",
        "Pipeline modes: Triggered vs Continuous, Dev vs Prod",
    ]),
    ("Lakeflow Jobs", ACCENT_PURPLE, [
        "DAG patterns: linear, fan-out/fan-in, conditional",
        "Repair runs: only failed + downstream re-executed",
        "4 trigger modes: manual, cron, continuous, file arrival",
        "DAB commands: bundle validate, deploy, run",
    ]),
]

for i, (title, color, items) in enumerate(sections):
    left = Inches(0.3 + i * 4.3)
    box = add_rect(slide, left, Inches(1.3), Inches(4.0), Inches(5.5), BG_SECTION, color)
    tf = set_text(box, title, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    tf.add_paragraph()
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(8)

# ============================================================
# SLIDE 25: Key Takeaways
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_BLUE)
text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "Key Takeaways", size=36, color=WHITE, bold=True)

takeaways = [
    ("1", "Lakeflow has 3 components: Connect (ingest), SDP (transform), Jobs (orchestrate)", ACCENT_ORANGE),
    ("2", "Connect offers Manual Upload, Standard Connectors (Auto Loader, JDBC, Kafka), and Managed Connectors", ACCENT_ORANGE),
    ("3", "SDP is declarative: define WHAT, framework handles sequencing, parallelism, retries, checkpoints", ACCENT_GREEN),
    ("4", "Three dataset types: Streaming Table (append), Materialized View (recompute), View (no storage)", ACCENT_GREEN),
    ("5", "Expectations enforce data quality: expect (monitor), expect_or_drop (filter), expect_or_fail (halt)", ACCENT_GREEN),
    ("6", "Auto CDC Flow replaces manual MERGE with declarative SCD Type 1 and Type 2", ACCENT_GREEN),
    ("7", "Jobs supports multi-task DAGs with repair runs (re-run only failed tasks)", ACCENT_PURPLE),
    ("8", "DAB brings CI/CD: jobs as code, automated deployment across dev/staging/prod", ACCENT_PURPLE),
]

for i, (num, text, color) in enumerate(takeaways):
    top = Inches(1.2 + i * 0.73)
    circle = add_rect(slide, Inches(0.5), top, Inches(0.5), Inches(0.5), color)
    set_text(circle, num, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, Inches(1.2), top, Inches(11.5), Inches(0.5),
             text, size=16, color=LIGHT_GRAY)

# ============================================================
# SLIDE 26: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_SECTION)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)
text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
         "Thank You!", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
         "Next: Day 25 -- SCD Type 2 Pipelines  |  Day 26 -- Performance Engineering", size=22, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.06), ACCENT_RED)
text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "Spark Databricks Zero-to-Pro", size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = "resources/lakeflow-connect-sdp-jobs.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
