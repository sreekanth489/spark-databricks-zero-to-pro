"""
Generate PowerPoint: Unity Catalog Deep Dive
Includes all provided images with explanations.
Run:  python3 resources/generate_uc_deep_dive_pptx.py
Out:  resources/unity-catalog-deep-dive.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette ──────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x1B, 0x1F, 0x23)
BG_CARD    = RGBColor(0x23, 0x29, 0x2E)
BG_GREEN   = RGBColor(0x1B, 0x5E, 0x20)
BG_RED     = RGBColor(0x7F, 0x0E, 0x0E)
ACCENT_RED = RGBColor(0xE8, 0x3E, 0x3E)
ACCENT_BLU = RGBColor(0x42, 0x9E, 0xF5)
ACCENT_GRN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_ORG = RGBColor(0xFF, 0x98, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY      = RGBColor(0xB0, 0xBE, 0xC5)
DGRAY      = RGBColor(0x60, 0x60, 0x60)
YELLOW     = RGBColor(0xFF, 0xD6, 0x00)

IMG_DIR = os.path.join(os.path.dirname(__file__), "images")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


# ── Helpers ────────────────────────────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=BG_DARK):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, border=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, size=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def bullets(tf, items, size=16, color=LGRAY, bullet="•"):
    for item in items:
        p = tf.add_paragraph()
        p.text = f"{bullet}  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0


def title_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, SW, Inches(1.1), BG_CARD, radius=False)
    tf = txt(slide, Inches(0.4), Inches(0.12), SW - Inches(0.8), Inches(0.55),
             title, size=32, bold=True, color=WHITE)
    if subtitle:
        txt(slide, Inches(0.4), Inches(0.68), SW - Inches(0.8), Inches(0.38),
            subtitle, size=17, color=LGRAY)


def img(slide, path, l, t, w, h):
    full = os.path.join(IMG_DIR, path)
    if os.path.exists(full):
        slide.shapes.add_picture(full, l, t, w, h)


def credit(slide, text="Credit: Databricks, Inc."):
    txt(slide, Inches(0.2), SH - Inches(0.35), SW - Inches(0.4), Inches(0.3),
        text, size=10, color=DGRAY, align=PP_ALIGN.RIGHT)


def divider(slide, y, color=ACCENT_BLU):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.4), y, SW - Inches(0.8), Pt(2))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
rect(s, 0, 0, SW, SH, BG_DARK, radius=False)

# accent bar
rect(s, 0, Inches(3.05), SW, Inches(0.06), ACCENT_BLU, radius=False)

txt(s, Inches(1), Inches(1.2), SW - Inches(2), Inches(1.0),
    "Unity Catalog", size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(2.3), SW - Inches(2), Inches(0.7),
    "From Hive Metastore to Enterprise Data Governance",
    size=26, color=ACCENT_BLU, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.2), SW - Inches(2), Inches(0.5),
    "Metastore  ·  Before & After  ·  3-Level Namespace  ·  "
    "Row Filters  ·  Column Masks  ·  Lineage  ·  Lakehouse Federation",
    size=15, color=LGRAY, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.8), SW - Inches(2), Inches(0.4),
    "Databricks Zero-to-Pro  |  Day 10 & 11  |  Data Governance Module",
    size=13, color=DGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Agenda", "What we will cover today")

topics = [
    ("1", "Why do we need a Metastore?",          "Concept: data files + metadata must be in sync"),
    ("2", "Apache Hive Metastore",                 "The original — what it does and what it lacks"),
    ("3", "Why Unity Catalog?",                    "The real problems HMS cannot solve"),
    ("4", "Before & After Unity Catalog",          "Architecture comparison — workspace isolation vs account-level"),
    ("5", "3-Level Namespace",                     "Metastore → Catalog → Schema → Table/View/Volume"),
    ("6", "Databricks Lakehouse Governance",       "One governance model for data + AI"),
    ("7", "Admin Roles Hierarchy",                 "Account Admin → Metastore Admin → Workspace Admin"),
    ("8", "Row Filters & Column Masks",            "Native table-level security — replaces regional views"),
    ("9", "Data Lineage",                          "Automatic lineage from files to dashboards"),
    ("10", "Lakehouse Federation",                 "Query Snowflake/MySQL/PostgreSQL live — no copy"),
]
col_w = Inches(6.2)
for i, (num, topic, sub) in enumerate(topics):
    row = i % 5
    col = i // 5
    x = Inches(0.3) + col * col_w
    y = Inches(1.25) + row * Inches(1.15)
    r = rect(s, x, y, col_w - Inches(0.15), Inches(1.0), BG_CARD,
             border=ACCENT_BLU)
    txt(s, x + Inches(0.12), y + Inches(0.08), Inches(0.4), Inches(0.4),
        num, size=20, bold=True, color=ACCENT_BLU)
    txt(s, x + Inches(0.55), y + Inches(0.08), col_w - Inches(0.75),
        Inches(0.38), topic, size=16, bold=True, color=WHITE)
    txt(s, x + Inches(0.55), y + Inches(0.52), col_w - Inches(0.75),
        Inches(0.4), sub, size=12, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What is a Metastore? (hand-drawn image)
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "What is a Metastore?",
          "A table = Data Files in cloud storage  +  Metadata in the Metastore — both must be in sync")

# image — left half
img(s, "metastore-concept.png", Inches(0.3), Inches(1.3), Inches(5.2), Inches(4.5))
credit(s, "Diagram: Sreekanth Keerthipati")

# explanation — right half
x = Inches(6.0)
txt(s, x, Inches(1.3), Inches(7.0), Inches(0.45),
    "Why does a Metastore exist?", size=20, bold=True, color=ACCENT_BLU)

tf = txt(s, x, Inches(1.85), Inches(7.0), Inches(0.35),
         "When you store data in S3 / ADLS, you get files. But to run SQL,",
         size=15, color=LGRAY)

points = [
    "WHERE is the file?   (s3://databricks/orders)",
    "WHAT is the schema? (id INT, name STRING, amount DOUBLE)",
    "WHAT format?        (Delta / Parquet / CSV)",
    "WHO can access it?  (ACLs / Row Filters / Column Masks)",
]
tf2 = txt(s, x, Inches(2.25), Inches(7.0), Inches(1.6),
          "", size=14, color=LGRAY)
bullets(tf2, points, size=14, color=LGRAY, bullet="→")

divider(s, Inches(4.05))

txt(s, x, Inches(4.15), Inches(7.0), Inches(0.4),
    "⚠  If Metastore & Storage are OUT OF SYNC:", size=15, bold=True,
    color=ACCENT_ORG)

problems = [
    "Files deleted from S3, metastore not updated  →  query fails",
    "Files added to S3 without metastore update  →  data invisible to SQL",
    "Table dropped, files left in S3  →  ghost files wasting storage cost",
]
tf3 = txt(s, x, Inches(4.6), Inches(7.0), Inches(1.6),
          "", size=13, color=LGRAY)
bullets(tf3, problems, size=13, color=ACCENT_ORG, bullet="✗")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Apache Hive Metastore
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Apache Hive Metastore (HMS)",
          "The original open-source metastore — used by Databricks before Unity Catalog")

# Left: What HMS does
xl = Inches(0.3)
rect(s, xl, Inches(1.25), Inches(6.2), Inches(5.8), BG_CARD,
     border=ACCENT_GRN)
txt(s, xl + Inches(0.15), Inches(1.35), Inches(5.9), Inches(0.4),
    "✓  What HMS provides", size=18, bold=True, color=ACCENT_GRN)
tf = txt(s, xl + Inches(0.15), Inches(1.85), Inches(5.9), Inches(0.3),
         "", size=14, color=LGRAY)
bullets(tf, [
    "Stores table DDL, schema, location, partitions, stats",
    "2-level namespace:  schema.table",
    "Open-source — compatible with Hive, Spark, Trino, Presto",
    "Backend: MySQL / PostgreSQL / Derby relational DB",
    "Supports tables, views, and functions",
    "GRANT / REVOKE on tables and databases",
], size=14, color=LGRAY)

# Right: What HMS lacks
xr = Inches(6.85)
rect(s, xr, Inches(1.25), Inches(6.2), Inches(5.8), BG_CARD,
     border=ACCENT_RED)
txt(s, xr + Inches(0.15), Inches(1.35), Inches(5.9), Inches(0.4),
    "✗  What HMS cannot do", size=18, bold=True, color=ACCENT_RED)
tf2 = txt(s, xr + Inches(0.15), Inches(1.85), Inches(5.9), Inches(0.3),
          "", size=14, color=LGRAY)
bullets(tf2, [
    "Workspace-isolated — no cross-workspace sharing",
    "Users / groups must be re-created in every workspace",
    "No governance for files (S3 IAM = all-or-nothing)",
    "No data lineage — impossible to trace data origin",
    "No audit logs — no trail of who accessed what",
    "No governance for ML models or AI assets",
    "Regional views required for row-level filtering",
    "Column masking only via complex dynamic views",
], size=14, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Why Unity Catalog? (Problems HMS cannot solve)
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Why Unity Catalog?",
          "The real-world problems that Hive Metastore could not solve — and Unity Catalog does")

problems_uc = [
    ("🏝  Workspace Islands",
     "Each workspace had its own HMS with separate users, groups, and ACLs.\n"
     "Moving between workspaces meant re-creating everything from scratch."),
    ("🔁  Duplicate User Management",
     "A team of 50 people across 5 workspaces = 250 user entries to maintain.\n"
     "When someone left, you had to remove them from every workspace manually."),
    ("🚫  No Cross-Workspace Sharing",
     "To share a table between Dev and Prod, you had to copy the data.\n"
     "Two copies = double storage cost + always slightly stale."),
    ("📂  Files Had No Governance",
     "Files in S3 / ADLS were protected only by IAM roles — all or nothing.\n"
     "UC Volumes bring the same GRANT/REVOKE model to files."),
    ("🕵️  No Audit Trail",
     "If someone accessed sensitive customer data at 2am, you had no idea.\n"
     "UC logs every access to system.access.audit — queryable with SQL."),
    ("🔗  No Data Lineage",
     "Teams manually maintained spreadsheets documenting data flows.\n"
     "UC captures lineage automatically — every query, job, dashboard."),
]

cols = 3
for i, (heading, body) in enumerate(problems_uc):
    row = i // cols
    col = i % cols
    card_w = Inches(4.1)
    card_h = Inches(2.4)
    x = Inches(0.2) + col * (card_w + Inches(0.15))
    y = Inches(1.3) + row * (card_h + Inches(0.15))
    rect(s, x, y, card_w, card_h, BG_CARD, border=ACCENT_RED)
    txt(s, x + Inches(0.15), y + Inches(0.12), card_w - Inches(0.3),
        Inches(0.42), heading, size=15, bold=True, color=ACCENT_RED)
    txt(s, x + Inches(0.15), y + Inches(0.58), card_w - Inches(0.3),
        card_h - Inches(0.7), body, size=12, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Before & After Unity Catalog (provided image)
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Before & After Unity Catalog",
          "Workspace-isolated Hive Metastores  →  Account-level shared governance")

img(s, "uc-before-after.png", Inches(0.5), Inches(1.25), Inches(12.3), Inches(5.3))
credit(s, "Credit: Databricks, Inc.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Before vs After Feature Table
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Hive Metastore vs Unity Catalog — Feature Comparison")

headers = ["Capability", "Hive Metastore", "Unity Catalog"]
rows = [
    ["Scope",              "Per workspace (isolated)",       "Account-level (all workspaces)"],
    ["Namespace",          "schema.table  (2-level)",        "catalog.schema.table  (3-level)"],
    ["User management",    "Duplicated per workspace",       "Single identity, assigned to workspaces"],
    ["File governance",    "None (IAM = all or nothing)",    "Volumes with GRANT/REVOKE"],
    ["Row-level security", "Dynamic views only",             "Native Row Filters on table"],
    ["Column masking",     "Dynamic views only",             "Native Column Masks on column"],
    ["Data lineage",       "None",                           "Automatic — tables, jobs, dashboards"],
    ["Audit logging",      "None",                           "Full — system.access.audit"],
    ["ML model governance","None",                           "Models as governed objects"],
    ["Cross-workspace",    "Data copy required",             "Native shared metastore"],
]

col_widths = [Inches(2.8), Inches(4.5), Inches(5.5)]
col_x = [Inches(0.2), Inches(3.1), Inches(7.7)]
header_y = Inches(1.2)
row_h = Inches(0.52)

# Header row
for ci, (h, w, x) in enumerate(zip(headers, col_widths, col_x)):
    fill = ACCENT_BLU if ci == 0 else (BG_RED if ci == 1 else BG_GREEN)
    rect(s, x, header_y, w, Inches(0.45), fill, radius=False)
    txt(s, x + Inches(0.08), header_y + Inches(0.05),
        w - Inches(0.1), Inches(0.38), h, size=14, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

for ri, row_data in enumerate(rows):
    y = header_y + Inches(0.48) + ri * row_h
    fill_row = RGBColor(0x1E, 0x24, 0x28) if ri % 2 == 0 else BG_CARD
    for ci, (cell, w, x) in enumerate(zip(row_data, col_widths, col_x)):
        rect(s, x, y, w, row_h - Inches(0.04), fill_row, radius=False)
        c = WHITE if ci == 0 else (ACCENT_RED if ci == 1 else ACCENT_GRN)
        txt(s, x + Inches(0.1), y + Inches(0.07),
            w - Inches(0.15), row_h - Inches(0.1),
            cell, size=12, color=c)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 3-Level Namespace (provided image)
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Unity Catalog — 3-Level Namespace",
          "Metastore → Catalog → Schema (Database) → Table / View / Volume / Function")

img(s, "uc-three-level-namespace.png",
    Inches(0.3), Inches(1.2), Inches(7.5), Inches(5.4))
credit(s, "Credit: Databricks, Inc.")

# Key points right side
xr = Inches(8.1)
txt(s, xr, Inches(1.3), Inches(5.0), Inches(0.45),
    "What changed?", size=20, bold=True, color=ACCENT_BLU)

txt(s, xr, Inches(1.85), Inches(5.0), Inches(0.35),
    "Hive Metastore:  schema.table  (2 levels)",
    size=15, color=ACCENT_RED, bold=True)
txt(s, xr, Inches(2.25), Inches(5.0), Inches(0.35),
    "Unity Catalog:   catalog.schema.table  (3 levels)",
    size=15, color=ACCENT_GRN, bold=True)

divider(s, Inches(2.75))

txt(s, xr, Inches(2.9), Inches(5.0), Inches(0.38),
    "Why a Catalog level?", size=17, bold=True, color=WHITE)
tf = txt(s, xr, Inches(3.35), Inches(5.0), Inches(0.3), "", size=14, color=LGRAY)
bullets(tf, [
    "Separate environments:  dev / staging / prod",
    "Separate business units:  finance / marketing / hr",
    "Separate projects or teams",
    "Different storage locations per catalog",
], size=14, color=LGRAY)

divider(s, Inches(5.3))
txt(s, xr, Inches(5.4), Inches(5.0), Inches(0.4),
    "Also under Metastore:", size=14, bold=True, color=ACCENT_ORG)
tf2 = txt(s, xr, Inches(5.85), Inches(5.0), Inches(0.3), "", size=13, color=LGRAY)
bullets(tf2, [
    "Storage Credential  (cloud auth)",
    "External Location   (S3/ADLS path)",
    "Share / Recipient   (Delta Sharing)",
], size=13, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Databricks Lakehouse Unifies Data & AI Governance
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Databricks Lakehouse Unifies Data & AI Governance",
          "Unity Catalog is the single governance layer — above: compute  |  below: storage & federation")

img(s, "databricks-lakehouse-unity-catalog.png",
    Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5))
credit(s, "Credit: Databricks, Inc.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Admin Roles Hierarchy (provided image)
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Databricks Account Roles Hierarchy",
          "Account Admin creates metastores · Metastore Admin governs data · Workspace Admin manages compute")

img(s, "databricks-roles-hierarchy.png",
    Inches(0.3), Inches(1.2), Inches(7.8), Inches(5.4))
credit(s, "Credit: Databricks, Inc.")

# Right: role descriptions
xr = Inches(8.4)
roles = [
    (ACCENT_RED, "Account Admin",
     ["Enable Unity Catalog, create metastores",
      "Create storage credentials",
      "Manage identities, integrate with Azure AD",
      "Create and manage workspaces",
      "Recommended: Platform Ops / Central Gov Team"]),
    (ACCENT_ORG, "Metastore Admin",
     ["Create / manage catalogs",
      "Grant privileges on any object",
      "Create storage credentials & external locations",
      "Transfer ownership of any object"]),
    (ACCENT_BLU, "Workspace Admin",
     ["Manage workspace users and groups",
      "Manage clusters, jobs, notebooks",
      "Cannot create metastores or storage creds"]),
]
y = Inches(1.3)
for color, role, points in roles:
    txt(s, xr, y, Inches(4.6), Inches(0.38), role,
        size=17, bold=True, color=color)
    tf = txt(s, xr, y + Inches(0.42), Inches(4.6), Inches(0.28),
             "", size=12, color=LGRAY)
    bullets(tf, points, size=12, color=LGRAY)
    y += Inches(0.42) + len(points) * Inches(0.3) + Inches(0.25)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Row Filters & Column Masks (concept)
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Row Filters & Column Masks",
          "Native table-level security — replaces the 'create a view per region' pattern")

# Old way (left)
rect(s, Inches(0.2), Inches(1.25), Inches(6.2), Inches(5.8),
     BG_CARD, border=ACCENT_RED)
txt(s, Inches(0.35), Inches(1.35), Inches(5.9), Inches(0.45),
    "✗  OLD: Regional Views Pattern", size=17, bold=True, color=ACCENT_RED)
old_items = [
    "CREATE VIEW apac_employees AS SELECT * WHERE region='APAC'",
    "CREATE VIEW emea_employees AS SELECT * WHERE region='EMEA'",
    "CREATE VIEW amer_employees AS SELECT * WHERE region='AMER'",
    "",
    "Problems:",
    "  • N regions = N views to maintain",
    "  • Users must know WHICH view to query",
    "  • View can be bypassed if user has TABLE access",
    "  • Column masks require ANOTHER set of views",
    "  • New region added → new view required",
]
tf = txt(s, Inches(0.35), Inches(1.9), Inches(5.9), Inches(0.28), "", size=13, color=LGRAY)
for item in old_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(12 if item.startswith("  •") else 13)
    p.font.color.rgb = ACCENT_RED if item.startswith("  •") or item == "Problems:" else LGRAY

# New way (right)
rect(s, Inches(6.85), Inches(1.25), Inches(6.2), Inches(5.8),
     BG_CARD, border=ACCENT_GRN)
txt(s, Inches(7.0), Inches(1.35), Inches(5.9), Inches(0.45),
    "✓  NEW: Native Row Filters + Column Masks", size=17, bold=True, color=ACCENT_GRN)

code_lines = [
    ("-- Create row filter function (ONCE)", DGRAY),
    ("CREATE FUNCTION region_filter(region STRING)", WHITE),
    ("RETURNS BOOLEAN", WHITE),
    ("RETURN is_account_group_member('admins')", WHITE),
    ("    OR is_account_group_member(lower(region));", WHITE),
    ("", WHITE),
    ("-- Attach to TABLE (ONCE)", DGRAY),
    ("ALTER TABLE employees", WHITE),
    ("  SET ROW FILTER region_filter ON (region);", WHITE),
    ("", WHITE),
    ("-- Create column mask (ONCE per column)", DGRAY),
    ("ALTER TABLE employees", WHITE),
    ("  ALTER COLUMN salary SET MASK mask_salary;", WHITE),
    ("", WHITE),
    ("-- Everyone queries the TABLE directly", DGRAY),
    ("SELECT * FROM employees;", ACCENT_GRN),
    ("-- Filter + masks applied AUTOMATICALLY", DGRAY),
]
yc = Inches(1.9)
for line, color in code_lines:
    txt(s, Inches(7.0), yc, Inches(5.9), Inches(0.28),
        line, size=11, color=color)
    yc += Inches(0.28)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Data Lineage (provided screenshot)
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Data Lineage in Unity Catalog",
          "Automatic lineage — from S3 files through bronze/silver/gold to materialized views")

img(s, "uc-data-lineage.png", Inches(0.3), Inches(1.2), Inches(9.5), Inches(5.5))
credit(s, "Screenshot: Sreekanth Keerthipati's Databricks workspace  |  ecommerce lakehouse pipeline")

xr = Inches(10.1)
txt(s, xr, Inches(1.3), Inches(3.0), Inches(0.45),
    "What this shows", size=18, bold=True, color=ACCENT_BLU)

tf = txt(s, xr, Inches(1.85), Inches(3.0), Inches(0.3), "", size=13, color=LGRAY)
bullets(tf, [
    "External S3 locations as sources",
    "→ ecommerce.bronze.orders",
    "→ ecommerce.silver.orders",
    "→ ecommerce.gold.fact_orders",
    "→ 4 regional materialized views",
    "  (west / southeast / northeast / midwest)",
    "",
    "Zero configuration needed —",
    "UC captures this automatically",
], size=12, color=LGRAY, bullet="")

divider(s, Inches(5.3))
txt(s, xr, Inches(5.4), Inches(3.0), Inches(0.38),
    "Use cases:", size=14, bold=True, color=ACCENT_ORG)
tf2 = txt(s, xr, Inches(5.85), Inches(3.0), Inches(0.28), "", size=12, color=LGRAY)
bullets(tf2, [
    "Impact analysis before schema changes",
    "Root cause when dashboard is wrong",
    "Compliance: trace PII flow",
], size=12, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Lakehouse Federation (provided image)
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Lakehouse Federation",
          "Discover, query, and govern all your data — no matter where it lives")

img(s, "lakehouse-federation.png", Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5))
credit(s, "Credit: Databricks, Inc.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Lakehouse Federation: How It Works
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Lakehouse Federation — How It Works",
          "Three UC objects: Connection → Foreign Catalog → Foreign Table")

# Three objects as cards
cards = [
    (ACCENT_BLU, "1  CONNECTION",
     "Stores credentials to reach an external system.\n"
     "Lives at the Metastore level.\n"
     "Created by Metastore Admin.",
     "CREATE CONNECTION snowflake_prod\n"
     "  TYPE SNOWFLAKE\n"
     "  OPTIONS (\n"
     "    host 'org.snowflake.com',\n"
     "    user 'databricks_user',\n"
     "    password secret('scope','key')\n"
     "  );"),
    (ACCENT_ORG, "2  FOREIGN CATALOG",
     "Maps a remote database into UC namespace.\n"
     "Schemas and tables are auto-discovered.\n"
     "Governed with same GRANT/REVOKE.",
     "CREATE FOREIGN CATALOG snowflake_cat\n"
     "  USING CONNECTION snowflake_prod\n"
     "  OPTIONS (database 'PROD_DB');\n"
     "\n"
     "SHOW SCHEMAS IN snowflake_cat;\n"
     "SHOW TABLES IN snowflake_cat.finance;"),
    (ACCENT_GRN, "3  FOREIGN TABLE",
     "Auto-discovered. Queryable via SQL.\n"
     "Reads live from source — no data copy.\n"
     "Pushdown: filters/aggregations in source.",
     "SELECT region, SUM(revenue)\n"
     "FROM snowflake_cat.finance.revenue\n"
     "WHERE year = 2024\n"
     "GROUP BY region;\n"
     "-- Runs IN Snowflake, returns result"),
]

card_w = Inches(4.1)
for i, (color, heading, desc, code) in enumerate(cards):
    x = Inches(0.2) + i * (card_w + Inches(0.2))
    rect(s, x, Inches(1.25), card_w, Inches(5.85), BG_CARD, border=color)
    txt(s, x + Inches(0.15), Inches(1.35), card_w - Inches(0.25),
        Inches(0.42), heading, size=17, bold=True, color=color)
    txt(s, x + Inches(0.15), Inches(1.85), card_w - Inches(0.25),
        Inches(1.1), desc, size=12, color=LGRAY)
    rect(s, x + Inches(0.12), Inches(3.1), card_w - Inches(0.25),
         Inches(3.7), RGBColor(0x0D, 0x11, 0x14), radius=False)
    txt(s, x + Inches(0.22), Inches(3.2), card_w - Inches(0.45),
        Inches(3.5), code, size=11, color=ACCENT_GRN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Summary & Key Takeaways
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
bg(s)
title_bar(s, "Key Takeaways", "Unity Catalog — what you need to remember")

takeaways = [
    (ACCENT_BLU, "Metastore = Metadata Registry",
     "Data files live in cloud storage (S3/ADLS). The metastore holds DDL, schema, location, ACLs. Both must stay in sync."),
    (ACCENT_RED, "HMS was workspace-local",
     "Hive Metastore gave each workspace its own isolated silo — no sharing, no cross-workspace governance, no lineage."),
    (ACCENT_GRN, "UC is account-level",
     "One metastore per region. Assign it to multiple workspaces. Define users/groups once, use everywhere."),
    (ACCENT_ORG, "3-Level Namespace",
     "catalog.schema.table — the catalog level organises by environment (dev/prod) or business unit."),
    (ACCENT_BLU, "Native Row Filters & Column Masks",
     "Attach a function to a TABLE — filter rows and mask columns transparently. No more N regional views."),
    (ACCENT_GRN, "Automatic Lineage",
     "Every SQL query, notebook, and job is tracked. See upstream sources and downstream consumers in Catalog Explorer."),
    (ACCENT_ORG, "Lakehouse Federation",
     "Query Snowflake, MySQL, PostgreSQL, Redshift live — no copy, no ETL. UC governs access with same GRANT/REVOKE."),
    (ACCENT_RED, "Admin Roles",
     "Account Admin → creates metastores. Metastore Admin → governs data. Workspace Admin → manages compute."),
]

cols = 2
card_w = Inches(6.3)
card_h = Inches(1.5)
for i, (color, heading, body) in enumerate(takeaways):
    col = i % cols
    row = i // cols
    x = Inches(0.2) + col * (card_w + Inches(0.25))
    y = Inches(1.25) + row * (card_h + Inches(0.12))
    rect(s, x, y, card_w, card_h, BG_CARD, border=color)
    rect(s, x, y, Inches(0.08), card_h, color, radius=False)
    txt(s, x + Inches(0.2), y + Inches(0.1), card_w - Inches(0.3),
        Inches(0.4), heading, size=15, bold=True, color=color)
    txt(s, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.3),
        card_h - Inches(0.6), body, size=12, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
OUT = os.path.join(os.path.dirname(__file__), "unity-catalog-deep-dive.pptx")
prs.save(OUT)
print(f"✓ Saved: {OUT}")
print(f"  Slides: {len(prs.slides)}")
