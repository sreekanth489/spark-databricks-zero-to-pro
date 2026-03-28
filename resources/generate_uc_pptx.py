"""
Generate PowerPoint presentation for Unity Catalog Fundamentals & Security.
Run: python3 resources/generate_uc_pptx.py
Output: resources/unity-catalog-fundamentals-and-security.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Color palette --
BG_DARK = RGBColor(0x1B, 0x1F, 0x23)
BG_SECTION = RGBColor(0x23, 0x29, 0x2E)
ACCENT_RED = RGBColor(0xE8, 0x3E, 0x3E)
ACCENT_BLUE = RGBColor(0x42, 0x9E, 0xF5)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
ACCENT_PURPLE = RGBColor(0xAB, 0x47, 0xBC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBE, 0xC5)
DARK_GRAY = RGBColor(0x60, 0x60, 0x60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_text_box(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_bullet_list(tf, items, size=16, color=LIGHT_GRAY, level=0):
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = level
        p.space_before = Pt(4)


def add_code_box(slide, left, top, width, height, code_text, font_size=14):
    shape = add_shape_rect(slide, left, top, width, height, RGBColor(0x12, 0x15, 0x18), ACCENT_BLUE)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(code_text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN
        p.font.name = "Courier New"
    return shape


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Unity Catalog", size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Fundamentals & Security", size=36, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
             "Centralized Governance for Data & AI on Databricks", size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(5.5), Inches(5.2), Inches(2.3), Inches(0.06), ACCENT_RED)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "Spark Databricks Zero-to-Pro  |  Day 10 & 11", size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Agenda", size=36, color=WHITE, bold=True)

items_left = [
    "Part 1: Unity Catalog Fundamentals",
    "  What is Unity Catalog?",
    "  3-Level Namespace",
    "  Metastore Architecture",
    "  Managed Storage",
    "  Volumes for File Governance",
    "  Data Lineage & Search",
    "  Legacy Hive Metastore Access",
]
items_right = [
    "Part 2: Unity Catalog Security",
    "  Security Model Overview",
    "  Identities & Identity Federation",
    "  Privileges & Ownership",
    "  Row-Level Security",
    "  Column Masking",
    "  Storage Credentials & External Locations",
    "  Best Practices",
]

tf = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
                  "", size=18, color=WHITE)
for item in items_left:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(18) if not item.startswith("  ") else Pt(16)
    p.font.color.rgb = ACCENT_ORANGE if not item.startswith("  ") else LIGHT_GRAY
    p.font.bold = not item.startswith("  ")
    p.space_before = Pt(6)

tf2 = add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5),
                   "", size=18, color=WHITE)
for item in items_right:
    p = tf2.add_paragraph()
    p.text = item
    p.font.size = Pt(18) if not item.startswith("  ") else Pt(16)
    p.font.color.rgb = ACCENT_PURPLE if not item.startswith("  ") else LIGHT_GRAY
    p.font.bold = not item.startswith("  ")
    p.space_before = Pt(6)

# ============================================================
# SLIDE 3: What is Unity Catalog?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "What is Unity Catalog?", size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
             "A centralized governance solution for all data and AI assets across every workspace and cloud.",
             size=20, color=ACCENT_BLUE, bold=True)

features = [
    "Centralized governance -- define access rules ONCE, enforce EVERYWHERE",
    "Unifies governance for tables, views, volumes, ML models, functions, dashboards",
    "Account-level identity management (not workspace-scoped)",
    "Automated lineage tracking across all asset types",
    "Built-in data search and discovery",
    "Full audit logging of all access and changes",
    "No hard migration -- legacy hive_metastore remains accessible",
]
tf = add_text_box(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5), "", size=18)
for feat in features:
    p = tf.add_paragraph()
    p.text = f"  {feat}"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

# ============================================================
# SLIDE 4: Before vs After Unity Catalog
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Before vs After Unity Catalog", size=36, color=WHITE, bold=True)

# Before box
add_shape_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(0x2D, 0x15, 0x15), ACCENT_RED)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.5),
             "Before: Hive Metastore", size=22, color=ACCENT_RED, bold=True)
before_items = [
    "Users & groups per workspace",
    "ACLs scoped to single workspace",
    "No file governance (ANY FILE only)",
    "No cross-workspace sharing",
    "No automated lineage",
    "No built-in search/discovery",
    "Tables & views only",
]
tf = add_text_box(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(4), "", size=16)
for item in before_items:
    p = tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(6)

# After box
add_shape_rect(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(0x15, 0x2D, 0x15), ACCENT_GREEN)
add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5.2), Inches(0.5),
             "After: Unity Catalog", size=22, color=ACCENT_GREEN, bold=True)
after_items = [
    "Account-level identity federation",
    "Cross-workspace, cross-cloud governance",
    "Volumes for governed file access",
    "Shared metastore across workspaces",
    "Automated lineage for all assets",
    "Built-in search and discovery",
    "Tables, views, volumes, models, functions",
]
tf = add_text_box(slide, Inches(7.3), Inches(2.3), Inches(5.2), Inches(4), "", size=16)
for item in after_items:
    p = tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(6)

# ============================================================
# SLIDE 5: 3-Level Namespace
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "The 3-Level Namespace", size=36, color=WHITE, bold=True)

# Hive vs UC
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5),
             "Hive Metastore (2-Level)", size=20, color=ACCENT_RED, bold=True)
add_code_box(slide, Inches(0.8), Inches(1.9), Inches(5), Inches(0.6),
             "SELECT * FROM schema.table", font_size=16)

add_text_box(slide, Inches(7), Inches(1.4), Inches(5), Inches(0.5),
             "Unity Catalog (3-Level)", size=20, color=ACCENT_GREEN, bold=True)
add_code_box(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(0.6),
             "SELECT * FROM catalog.schema.table", font_size=16)

# Hierarchy diagram
hierarchy_text = """Metastore (top-level container)
  |
  +-- Catalog: prod_catalog
  |     +-- Schema: hr_db
  |     |     +-- Table: employees
  |     |     +-- View: active_employees_vw
  |     |     +-- Volume: raw_files
  |     |     +-- Function: mask_ssn()
  |     +-- Schema: finance_db
  |           +-- Table: transactions
  |
  +-- Catalog: dev_catalog
  |     +-- Schema: sandbox
  |
  +-- Catalog: hive_metastore (legacy)
        +-- Schema: default"""
add_code_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(4.2), hierarchy_text, font_size=14)

# ============================================================
# SLIDE 6: Metastore Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_PURPLE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Metastore Architecture", size=36, color=WHITE, bold=True)

facts = [
    ("Regional", "One metastore per cloud region (e.g., us-east-1)"),
    ("Multi-Workspace", "A single metastore assigned to multiple workspaces"),
    ("Shared View", "All assigned workspaces see the same catalogs & tables"),
    ("Managed Storage", "S3/ADLS/GCS bucket for managed table data"),
    ("Metadata + ACLs", "Stores object definitions and access control lists"),
    ("Audit Trail", "Full log of who accessed what, when"),
]
colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_PURPLE, ACCENT_RED, ACCENT_BLUE]

for i, (title, desc) in enumerate(facts):
    row = i // 3
    col_idx = i % 3
    left = Inches(0.5 + col_idx * 4.2)
    top = Inches(1.5 + row * 2.8)
    box = add_shape_rect(slide, left, top, Inches(3.8), Inches(2.2), BG_SECTION, colors[i])
    tf = set_text(box, title, size=20, color=colors[i], bold=True, alignment=PP_ALIGN.CENTER)
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(12)

# ============================================================
# SLIDE 7: Managed Storage
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Managed Storage Hierarchy", size=36, color=WHITE, bold=True)

storage_text = """Metastore Default Storage: s3://uc-metastore-bucket/
  |
  +-- Catalog-level override (optional)
  |     s3://prod-catalog-bucket/
  |
  +-- Schema-level override (optional)
        s3://hr-schema-bucket/"""
add_code_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(2.2), storage_text, font_size=16)

# Managed vs External
add_shape_rect(slide, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.8), BG_SECTION, ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.5),
             "Managed Tables", size=22, color=ACCENT_GREEN, bold=True)
tf = add_text_box(slide, Inches(0.8), Inches(4.9), Inches(5), Inches(1.8), "", size=15)
add_bullet_list(tf, [
    "UC controls storage location",
    "DROP TABLE deletes data",
    "Users never see physical paths",
    "Default behavior in UC",
], size=15)

add_shape_rect(slide, Inches(7), Inches(4.2), Inches(5.8), Inches(2.8), BG_SECTION, ACCENT_ORANGE)
add_text_box(slide, Inches(7.3), Inches(4.3), Inches(5), Inches(0.5),
             "External Tables", size=22, color=ACCENT_ORANGE, bold=True)
tf = add_text_box(slide, Inches(7.3), Inches(4.9), Inches(5), Inches(1.8), "", size=15)
add_bullet_list(tf, [
    "You specify LOCATION to existing files",
    "DROP TABLE keeps data intact",
    "UC manages metadata & access only",
    "Use for shared/existing data lakes",
], size=15)

# ============================================================
# SLIDE 8: Volumes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Unity Catalog Volumes", size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "Governed access to non-tabular files: CSVs, images, JARs, model artifacts",
             size=20, color=ACCENT_BLUE, bold=True)

vol_code = """-- Create a managed volume
CREATE VOLUME my_catalog.my_schema.raw_files;

-- List files in a volume
LIST '/Volumes/my_catalog/my_schema/raw_files/';

-- Read files directly from a volume
SELECT * FROM read_files(
  '/Volumes/my_catalog/my_schema/raw_files/data.csv',
  format => 'csv', header => 'true'
);"""
add_code_box(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.5), vol_code, font_size=15)

tf = add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(1), "", size=16)
add_bullet_list(tf, [
    "Same GRANT/REVOKE permission model as tables (READ VOLUME, WRITE VOLUME)",
    "Managed volumes = UC controls storage | External volumes = you point to existing path",
], size=16)

# ============================================================
# SLIDE 9: Lineage & Search
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_PURPLE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Automated Lineage & Data Search", size=36, color=WHITE, bold=True)

lineage_text = """Source Table A ---+
                  |---> Silver Table ---> Gold Table ---> Dashboard
Source Table B ---+         |
                            v
                       ML Model"""
add_code_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(2.2), lineage_text, font_size=16)

features = [
    ("Automated Lineage", "Tracks origin & flow across tables, views, notebooks, jobs, dashboards, ML models"),
    ("No Configuration", "Lineage captured automatically -- zero setup required"),
    ("Data Search", "Discover tables, columns, descriptions through built-in search"),
    ("Audit Logging", "Full trail of who accessed what, when -- compliance-ready"),
]
for i, (title, desc) in enumerate(features):
    top = Inches(4.0 + i * 0.8)
    add_text_box(slide, Inches(1.5), top, Inches(2.5), Inches(0.6),
                 title, size=18, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(4.2), top, Inches(8), Inches(0.6),
                 desc, size=16, color=LIGHT_GRAY)

# ============================================================
# SLIDE 10: Section Divider - Security
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_SECTION)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
             "Part 2: Unity Catalog Security", size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             "RBAC, Privileges, Row-Level Security, Column Masking", size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_shape_rect(slide, Inches(5.5), Inches(5.0), Inches(2.3), Inches(0.06), ACCENT_RED)

# ============================================================
# SLIDE 11: Security Model Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "UC Security Model", size=36, color=WHITE, bold=True)

model_code = """GRANT  privilege  ON  securable_object  TO  principal

Principals (WHO):           Securable Objects (WHERE):     Privileges (WHAT):
+-------------------+      +------------------------+     +------------------+
| Users             |      | METASTORE              |     | SELECT           |
| Service Principals|      |   CATALOG              |     | MODIFY           |
| Groups            |      |     SCHEMA             |     | CREATE           |
+-------------------+      |       TABLE / VIEW     |     | USE CATALOG      |
                           |       VOLUME           |     | USE SCHEMA       |
                           |       FUNCTION / MODEL |     | READ/WRITE FILES |
                           | STORAGE CREDENTIAL     |     | READ/WRITE VOLUME|
                           | EXTERNAL LOCATION      |     | EXECUTE          |
                           | SHARE / RECIPIENT      |     | ALL PRIVILEGES   |
                           +------------------------+     +------------------+"""
add_code_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5), model_code, font_size=14)

# ============================================================
# SLIDE 12: Identities
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Identities (Principals)", size=36, color=WHITE, bold=True)

identities = [
    ("Users", ACCENT_BLUE, "Individual people identified by email address.\nCan be assigned admin roles."),
    ("Service Principals", ACCENT_GREEN, "Automated identities for tools & CI/CD.\nIdentified by Application ID."),
    ("Groups", ACCENT_ORANGE, "Collections of users & service principals.\nCan be nested (groups within groups)."),
]

for i, (title, color, desc) in enumerate(identities):
    left = Inches(0.5 + i * 4.2)
    box = add_shape_rect(slide, left, Inches(1.5), Inches(3.8), Inches(2.5), BG_SECTION, color)
    tf = set_text(box, title, size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)

# Identity Federation
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
             "Identity Federation", size=24, color=ACCENT_PURPLE, bold=True)
tf = add_text_box(slide, Inches(0.8), Inches(5.1), Inches(11), Inches(2), "", size=16)
add_bullet_list(tf, [
    "Identities created ONCE at account level, then assigned to workspaces as needed",
    "No duplicate identity management -- single source of truth",
    "Managed through the Account Console (accounts.cloud.databricks.com)",
], size=16)

# ============================================================
# SLIDE 13: Privilege Hierarchy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Privilege Hierarchy & Prerequisites", size=36, color=WHITE, bold=True)

prereq_code = """To SELECT from catalog.schema.table:

  Step 1: GRANT USE CATALOG ON CATALOG prod_catalog TO analysts;
  Step 2: GRANT USE SCHEMA ON SCHEMA prod_catalog.hr_db TO analysts;
  Step 3: GRANT SELECT ON TABLE prod_catalog.hr_db.employees TO analysts;

  USE CATALOG + USE SCHEMA = navigation rights (no data access)
  SELECT = actual read access to the data"""
add_code_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3), prereq_code, font_size=15)

# Owner hierarchy
roles = [
    ("Metastore Admin", "All objects in metastore"),
    ("Catalog Owner", "All objects in that catalog"),
    ("Schema Owner", "All objects in that schema"),
    ("Table Owner", "That specific table only"),
]
for i, (role, scope) in enumerate(roles):
    left = Inches(0.8 + i * 0.4)
    top = Inches(4.8 + i * 0.6)
    add_text_box(slide, left, top, Inches(3.5), Inches(0.5),
                 role, size=18, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(5), top, Inches(7), Inches(0.5),
                 f"Can grant on: {scope}", size=16, color=LIGHT_GRAY)

# ============================================================
# SLIDE 14: GRANT/REVOKE Syntax
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "GRANT / REVOKE Syntax", size=36, color=WHITE, bold=True)

grant_code = """-- Grant read access to a group
GRANT USE CATALOG ON CATALOG prod_catalog TO analysts;
GRANT USE SCHEMA ON SCHEMA prod_catalog.hr_db TO analysts;
GRANT SELECT ON TABLE prod_catalog.hr_db.employees TO analysts;

-- Grant multiple privileges
GRANT SELECT, MODIFY ON SCHEMA hr_db TO data_engineers;

-- Grant to individual user
GRANT SELECT ON VIEW hr_db.active_vw TO `analyst@company.com`;

-- Revoke privileges
REVOKE SELECT ON TABLE employees FROM analysts;
REVOKE ALL PRIVILEGES ON TABLE employees FROM data_engineers;

-- View grants
SHOW GRANTS ON TABLE employees;
SHOW GRANTS ON SCHEMA hr_db;
SHOW GRANTS `analyst@company.com`;"""
add_code_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8), grant_code, font_size=14)

# ============================================================
# SLIDE 15: Row-Level Security
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Row-Level Security with Dynamic Views", size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Filter rows based on the calling user's group membership", size=18, color=ACCENT_BLUE)

rls_code = """CREATE OR REPLACE VIEW secure_employees_vw AS
SELECT employee_id, first_name, last_name, email, department, salary
FROM employees
WHERE
  CASE
    WHEN is_account_group_member('admins') THEN true
    WHEN is_account_group_member('engineering')
         AND department = 'Engineering' THEN true
    WHEN is_account_group_member('marketing')
         AND department = 'Marketing' THEN true
    WHEN is_account_group_member('finance')
         AND department = 'Finance' THEN true
    ELSE false
  END;

-- Grant on VIEW (not the underlying table)
GRANT SELECT ON VIEW secure_employees_vw TO analysts;"""
add_code_box(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.2), rls_code, font_size=14)

# ============================================================
# SLIDE 16: Column Masking
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_PURPLE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Column Masking with Dynamic Views", size=36, color=WHITE, bold=True)

mask_code = """CREATE OR REPLACE VIEW masked_employees_vw AS
SELECT
  employee_id, first_name, last_name,

  -- Email: visible to HR/admins, masked for others
  CASE WHEN is_account_group_member('hr')
            OR is_account_group_member('admins')
       THEN email
       ELSE concat(left(email, 2), '***@***')
  END AS email,

  -- SSN: last 4 digits only for HR, fully masked for others
  CASE WHEN is_account_group_member('hr')
       THEN ssn
       ELSE concat('***-**-', right(ssn, 4))
  END AS ssn,

  -- Salary: visible to finance/HR/admins only
  CASE WHEN is_account_group_member('finance')
            OR is_account_group_member('hr')
       THEN salary  ELSE NULL
  END AS salary
FROM employees;"""
add_code_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8), mask_code, font_size=13)

# ============================================================
# SLIDE 17: Column Masking Matrix
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_PURPLE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Column Visibility Matrix", size=36, color=WHITE, bold=True)

# Table header
headers = ["Column", "HR Group", "Finance Group", "Engineering", "Admin"]
col_widths = [Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)]
start_left = Inches(0.8)
header_top = Inches(1.8)

for i, header in enumerate(headers):
    left = Inches(0.8 + i * 2.4)
    box = add_shape_rect(slide, left, header_top, Inches(2.2), Inches(0.6), ACCENT_BLUE)
    set_text(box, header, size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

rows = [
    ["email", "Full", "Masked", "Masked", "Full"],
    ["ssn", "Full", "Masked", "Masked", "Full"],
    ["salary", "Full", "Full", "NULL", "Full"],
    ["first_name", "Full", "Full", "Full", "Full"],
    ["department", "Full", "Full", "Full", "Full"],
]

for r, row in enumerate(rows):
    row_top = Inches(2.5 + r * 0.6)
    for c, cell in enumerate(row):
        left = Inches(0.8 + c * 2.4)
        if cell == "Masked" or cell == "NULL":
            cell_color = RGBColor(0x40, 0x20, 0x20)
            text_color = ACCENT_RED
        elif cell == "Full":
            cell_color = RGBColor(0x20, 0x40, 0x20)
            text_color = ACCENT_GREEN
        else:
            cell_color = BG_SECTION
            text_color = WHITE
        box = add_shape_rect(slide, left, row_top, Inches(2.2), Inches(0.5), cell_color)
        set_text(box, cell, size=15, color=text_color, bold=(c == 0), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
             "Pattern: Grant SELECT on the VIEW, not on the underlying TABLE",
             size=18, color=ACCENT_ORANGE, bold=True)

# ============================================================
# SLIDE 18: Storage Credentials & External Locations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Storage Credentials & External Locations", size=36, color=WHITE, bold=True)

storage_code = """-- Storage Credential: authentication to cloud storage
CREATE STORAGE CREDENTIAL my_s3_cred
WITH (AWS_IAM_ROLE = 'arn:aws:iam::123456789:role/uc-role');

-- External Location: maps credential to a specific path
CREATE EXTERNAL LOCATION my_raw_data
URL 's3://my-bucket/raw-data/'
WITH (STORAGE CREDENTIAL my_s3_cred);

-- Grant file access through external locations
GRANT READ FILES ON EXTERNAL LOCATION my_raw_data TO data_engineers;
GRANT WRITE FILES ON EXTERNAL LOCATION my_raw_data TO data_engineers;"""
add_code_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.8), storage_code, font_size=14)

hier_text = """Storage Credential (IAM Role / Managed Identity / Service Account)
  +-- External Location 1: s3://bucket/raw/
  +-- External Location 2: s3://bucket/curated/
  +-- External Location 3: s3://bucket/archive/"""
add_code_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(1.5), hier_text, font_size=14)

# ============================================================
# SLIDE 19: Hive vs UC Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Security: Hive Metastore vs Unity Catalog", size=36, color=WHITE, bold=True)

comparison_rows = [
    ("Aspect", "Hive Metastore", "Unity Catalog"),
    ("Identity scope", "Workspace", "Account"),
    ("Identity types", "Users, groups", "Users, SPs, groups"),
    ("Securable objects", "Catalog, schema, table, view", "All + volumes, models, locations"),
    ("Prerequisites", "USAGE", "USE CATALOG + USE SCHEMA"),
    ("File governance", "ANY FILE (all or nothing)", "Per external location"),
    ("Row security", "Dynamic views", "Dynamic views + group functions"),
    ("Audit logging", "Limited", "Full audit trail"),
    ("Cross-workspace", "Not possible", "Built-in"),
]

for r, (aspect, hive, uc) in enumerate(comparison_rows):
    top = Inches(1.3 + r * 0.62)
    is_header = (r == 0)
    bg_color = ACCENT_BLUE if is_header else BG_SECTION
    text_color = WHITE if is_header else LIGHT_GRAY
    bold = is_header

    box1 = add_shape_rect(slide, Inches(0.5), top, Inches(3.2), Inches(0.55), bg_color)
    set_text(box1, aspect, size=14, color=text_color, bold=bold, alignment=PP_ALIGN.CENTER)

    box2 = add_shape_rect(slide, Inches(3.8), top, Inches(4.4), Inches(0.55), bg_color)
    set_text(box2, hive, size=14, color=ACCENT_RED if not is_header else text_color, bold=bold, alignment=PP_ALIGN.CENTER)

    box3 = add_shape_rect(slide, Inches(8.3), top, Inches(4.5), Inches(0.55), bg_color)
    set_text(box3, uc, size=14, color=ACCENT_GREEN if not is_header else text_color, bold=bold, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 20: Best Practices
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Security Best Practices", size=36, color=WHITE, bold=True)

practices = [
    ("Use groups, not individuals", "Assign privileges to groups for maintainability and lifecycle"),
    ("Least privilege principle", "Grant only the minimum privileges needed for the task"),
    ("Views for security", "Use dynamic views for row/column security -- never expose base tables"),
    ("Group ownership", "Assign table ownership to groups, not individual users"),
    ("Separate catalogs", "dev_catalog / staging_catalog / prod_catalog for environment isolation"),
    ("External locations", "Govern cloud storage through UC, not raw IAM alone"),
    ("Regular audits", "Review SHOW GRANTS and system audit logs periodically"),
]

for i, (title, desc) in enumerate(practices):
    top = Inches(1.3 + i * 0.82)
    add_shape_rect(slide, Inches(0.5), top, Inches(0.08), Inches(0.65), ACCENT_GREEN)
    add_text_box(slide, Inches(0.8), top, Inches(3.5), Inches(0.65),
                 title, size=17, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(4.5), top, Inches(8), Inches(0.65),
                 desc, size=16, color=LIGHT_GRAY)

# ============================================================
# SLIDE 21: Certification Tips
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Certification Tips", size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Databricks Certified Data Engineer Associate", size=18, color=ACCENT_BLUE, bold=True)

cert_items = [
    "3-level namespace: catalog.schema.table",
    "Difference between UC metastore and Hive Metastore",
    "Managed vs external tables -- lifecycle and storage behavior",
    "GRANT/REVOKE syntax and behavior",
    "USE CATALOG and USE SCHEMA as prerequisite privileges",
    "Ownership and who can grant (metastore admin > catalog owner > schema owner)",
    "is_account_group_member() for dynamic views",
    "Identity federation: account-level vs workspace-level",
    "Storage credentials and external locations",
    "Volumes for governed file access",
]

tf = add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(5), "", size=16)
for item in cert_items:
    p = tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(17)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

# ============================================================
# SLIDE 22: Key Takeaways
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Key Takeaways", size=36, color=WHITE, bold=True)

takeaways = [
    ("1", "Unity Catalog replaces workspace-scoped Hive Metastore with account-level governance"),
    ("2", "3-level namespace (catalog.schema.table) adds catalogs as the top organizational layer"),
    ("3", "Metastores are regional and can be shared across multiple workspaces"),
    ("4", "GRANT/REVOKE with USE CATALOG + USE SCHEMA prerequisites control access"),
    ("5", "Dynamic views implement row-level security and column masking"),
    ("6", "Storage credentials + external locations govern cloud storage access"),
    ("7", "Volumes extend governance to non-tabular files (CSVs, images, JARs)"),
    ("8", "Legacy hive_metastore remains accessible -- no hard migration required"),
]

for i, (num, text) in enumerate(takeaways):
    top = Inches(1.3 + i * 0.73)
    circle = add_shape_rect(slide, Inches(0.8), top, Inches(0.5), Inches(0.5), ACCENT_BLUE)
    set_text(circle, num, size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), top, Inches(11), Inches(0.5),
                 text, size=17, color=LIGHT_GRAY)

# ============================================================
# SLIDE 23: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_SECTION)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_RED)
add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
             "Thank You!", size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
             "Next: Day 12 -- Managed vs External Tables", size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_shape_rect(slide, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.06), ACCENT_RED)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "Spark Databricks Zero-to-Pro", size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = "resources/unity-catalog-fundamentals-and-security.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
